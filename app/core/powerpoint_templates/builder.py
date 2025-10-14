from pptx import Presentation
from pptx.util import Inches, Pt
import io
import logging


logger = logging.getLogger(__name__)


class TemplatePresentationBuilder:
    def __init__(self, template_path: str = None):
        if template_path:
            self.prs = Presentation(template_path)
            self._clear_existing_slides()  # ← ОЧИЩАЕМ СУЩЕСТВУЮЩИЕ СЛАЙДЫ
            logger.info(f"Загружен и очищен шаблон: {template_path}")
        else:
            self.prs = Presentation()
            self._setup_default_layout()

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def _clear_existing_slides(self):
        """Очищает все существующие слайды из шаблона"""
        try:
            # Создаем копию списка слайдов чтобы избежать проблем с итерацией
            slide_ids = [slide.slide_id for slide in self.prs.slides]

            for slide_id in slide_ids:
                # Находим слайд по ID и удаляем его
                for i, slide in enumerate(self.prs.slides):
                    if slide.slide_id == slide_id:
                        # Получаем r_id слайда
                        r_id = self.prs.slides._sldIdLst[i].rId
                        # Удаляем связь
                        self.prs.part.drop_rel(r_id)
                        break

            # Очищаем список слайдов
            self.prs.slides._sldIdLst = []

            logger.info("Все существующие слайды из шаблона удалены")
        except Exception as e:
            logger.warning(f"Не удалось очистить слайды из шаблона: {e}")

    def _setup_default_layout(self):
        """Настраивает базовые стили если шаблон не загружен"""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        logger.info("Используется стандартный макет презентации")

    def add_slide_from_template(self, slide_type: str, title: str, content: str,
                                layout_index: int = 1, audience: str = "инвесторы"):
        """Добавляет слайд используя макет из шаблона"""

        try:
            # Выбираем макет
            if slide_type == "title":
                layout_idx = 0  # Титульный макет
            else:
                # Используем контентные макеты по кругу
                available_layouts = len(self.prs.slide_layouts)
                if available_layouts > 1:
                    layout_idx = 1 + ((layout_index - 1) % (available_layouts - 1))
                else:
                    layout_idx = 0

            slide_layout = self.prs.slide_layouts[layout_idx]
            slide = self.prs.slides.add_slide(slide_layout)

            logger.info(f"Создан слайд с макетом {layout_idx}: {slide_layout.name}")

            # Заполняем плейсхолдеры
            success = self._fill_placeholders(slide, title, content, slide_type, audience)

            if not success:
                logger.warning(f"Не удалось заполнить плейсхолдеры для слайда: {title}")
                # Пробуем альтернативный метод
                self._fill_placeholders_fallback(slide, title, content)

            return slide

        except Exception as e:
            logger.error(f"Ошибка создания слайда из шаблона: {e}")
            return self._add_fallback_slide(title, content, slide_type)

    def _fill_placeholders(self, slide, title: str, content: str, slide_type: str, audience: str) -> bool:
        """Заполняет плейсхолдеры в слайде - возвращает успех"""
        try:
            title_filled = False
            content_filled = False

            # СПОСОБ 1: Используем стандартные плейсхолдеры
            if slide.shapes.title:
                slide.shapes.title.text = title
                title_filled = True
                logger.debug("Заголовок заполнен через shapes.title")

            # Ищем плейсхолдер для контента
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                if content_placeholder.has_text_frame:
                    self._format_content(content_placeholder.text_frame, content)
                    content_filled = True
                    logger.debug("Контент заполнен через placeholders[1]")

            # СПОСОБ 2: Ищем по типам фигур
            if not title_filled or not content_filled:
                for shape in slide.shapes:
                    if not title_filled and shape.has_text_frame:
                        # Пытаемся определить заголовок по тексту или позиции
                        if (shape.text and len(shape.text.strip()) < 50) or shape.top < Inches(1):
                            if "click to add title" in shape.text.lower() or "заголовок" in shape.text.lower():
                                shape.text = title
                                title_filled = True
                                logger.debug("Заголовок заполнен через поиск по тексту")

                    if not content_filled and shape.has_text_frame:
                        # Пытаемся определить контент по размеру или позиции
                        if (shape.height > Inches(2) and shape.top > Inches(1.5) and
                                ("click to add text" in shape.text.lower() or "текст" in shape.text.lower() or
                                 shape == content_placeholder)):
                            self._format_content(shape.text_frame, content)
                            content_filled = True
                            logger.debug("Контент заполнен через поиск по размеру")

            logger.info(f"Результат заполнения: заголовок={title_filled}, контент={content_filled}")
            return title_filled and content_filled

        except Exception as e:
            logger.warning(f"Ошибка в основном методе заполнения: {e}")
            return False

    def _fill_placeholders_fallback(self, slide, title: str, content: str):
        """Альтернативный метод заполнения плейсхолдеров"""
        try:
            logger.info("Используем альтернативный метод заполнения")

            # Просто ищем все текстовые фигуры и заполняем их
            text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]

            if len(text_shapes) >= 1:
                # Первая фигура - заголовок
                text_shapes[0].text = title

            if len(text_shapes) >= 2:
                # Вторая фигура - контент
                self._format_content(text_shapes[1].text_frame, content)
            elif len(text_shapes) == 1:
                # Если только одна фигура - используем ее для всего
                self._format_content(text_shapes[0].text_frame, f"{title}\n\n{content}")

        except Exception as e:
            logger.error(f"Ошибка в альтернативном методе: {e}")

    @staticmethod
    def _format_content(text_frame, content: str):
        """Форматирует контент в текстовом фрейме"""
        try:
            text_frame.clear()

            # Разбиваем контент на абзацы
            paragraphs = [p for p in content.split('\n') if p.strip()]

            for i, para in enumerate(paragraphs):
                p = text_frame.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(16)
                p.space_after = Pt(8)

                # Первый абзац без отступа, остальные с отступом
                if i > 0:
                    p.level = 1

            logger.debug(f"Отформатировано {len(paragraphs)} абзацев контента")

        except Exception as e:
            logger.error(f"Ошибка форматирования контента: {e}")
            # Простой fallback
            text_frame.text = content

    def _add_fallback_slide(self, title: str, content: str, slide_type: str):
        """Создает слайд по умолчанию если шаблон не сработал"""
        try:
            from app.core.pptx_builder import PresentationBuilder
            fallback_builder = PresentationBuilder()

            if slide_type == "title":
                slide = fallback_builder.add_title_slide(title, content)
            else:
                slide = fallback_builder.add_content_slide(title, content)

            logger.info("Использован fallback слайд")
            return slide

        except Exception as e:
            logger.error(f"Ошибка создания fallback слайда: {e}")
            # Создаем абсолютно простой слайд
            slide_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title
            return slide

    def save_to_bytes(self) -> io.BytesIO:
        """Сохраняет презентацию в bytes - надежная версия"""
        try:
            presentation_bytes = io.BytesIO()

            # Проверяем что есть слайды
            if len(self.prs.slides) == 0:
                logger.warning("Презентация пустая, добавляем заглушку")
                self.add_slide_from_template("title", "Пустая презентация", "Добавьте контент")

            # Сохраняем
            self.prs.save(presentation_bytes)
            presentation_bytes.seek(0)

            # Проверяем размер файла
            file_size = presentation_bytes.getbuffer().nbytes
            if file_size < 1000:  # Минимальный размер для PPTX
                raise ValueError(f"Файл слишком мал: {file_size} байт")

            logger.info(f"Презентация сохранена. Размер: {file_size} байт, слайдов: {len(self.prs.slides)}")
            return presentation_bytes

        except Exception as e:
            logger.error(f"Ошибка сохранения презентации с шаблоном: {e}")
            # Используем стандартный билдер как запасной вариант
            return self._create_fallback_simple_presentation()

    @staticmethod
    def _create_fallback_simple_presentation() -> io.BytesIO:
        """Создает простую презентацию без шаблона"""
        try:
            from app.core.pptx_builder import PresentationBuilder
            fallback_builder = PresentationBuilder()
            fallback_builder.add_title_slide(
                "Презентация AI Assistant",
                "Использован стандартный шаблон (исходный не сработал)"
            )
            return fallback_builder.save_to_bytes()
        except Exception as e:
            logger.error(f"Ошибка создания fallback презентации: {e}")
            raise

    def get_slide_count(self) -> int:
        """Возвращает количество слайдов"""
        return len(self.prs.slides)

    def get_available_layouts(self) -> list:
        """Возвращает список доступных макетов"""
        return [layout.name for layout in self.prs.slide_layouts]