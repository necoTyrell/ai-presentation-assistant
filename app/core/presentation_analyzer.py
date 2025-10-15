from pptx import Presentation
from typing import List, Dict
from pydantic import BaseModel
import logging

logger = logging.getLogger(__name__)


class TemplateAnalysis(BaseModel):
    slides_count: int
    layouts: List[Dict]
    master_slides: List[str]


def analyze_template(template_path: str) -> TemplateAnalysis:
    """Анализирует PowerPoint шаблон и возвращает информацию о макетах"""
    try:
        prs = Presentation(template_path)

        layouts = []

        # Анализируем доступные макеты
        for i, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "index": i,
                "name": layout.name,
                "placeholders": []
            }

            # Анализируем плейсхолдеры в макете
            for placeholder in layout.placeholders:
                if placeholder.has_text_frame:
                    layout_info["placeholders"].append({
                        "type": "text",
                        "index": placeholder.placeholder_format.idx,
                        "name": f"Placeholder {placeholder.placeholder_format.idx}"
                    })

            layouts.append(layout_info)

        # Получаем мастер-слайды
        master_slides = [master.name for master in prs.slide_masters if master.name]

        return TemplateAnalysis(
            slides_count=len(prs.slides),
            layouts=layouts,
            master_slides=master_slides
        )

    except Exception as e:
        logger.error(f"Ошибка анализа шаблона: {e}")
        raise