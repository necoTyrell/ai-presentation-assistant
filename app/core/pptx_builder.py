from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import logging

logger = logging.getLogger(__name__)


class PresentationBuilder:
    def __init__(self, template_path: str = None):
        if template_path:
            self.prs = Presentation(template_path)
            self._clear_slides()
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

    def _clear_slides(self):
        """Очищает все слайды из шаблона - ПРОСТОЙ ВАРИАНТ"""
        try:
            # Просто удаляем все существующие слайды
            slide_ids = [slide.slide_id for slide in self.prs.slides]
            for slide_id in slide_ids:
                self.prs.slides._sldIdLst.remove(slide_id)
            logger.info("✅ Шаблон очищен")
        except Exception as e:
            logger.warning(f"Не удалось очистить шаблон: {e}")
            # Создаем новую презентацию если не получилось очистить
            self.prs = Presentation()

    def add_slide(self, slide_type: str, title: str, content: str):
        layout_idx = 0 if slide_type == "title" else 1

        try:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        except:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        # Заголовок
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Контент с улучшенным форматированием
        try:
            if len(slide.placeholders) > 1:
                placeholder = slide.placeholders[1]
                text_frame = placeholder.text_frame
                text_frame.text = content

                # Форматируем текст
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.alignment = PP_ALIGN.LEFT

            else:
                # Создаем текстовое поле с правильными размерами
                left = Inches(1)
                top = Inches(2)
                width = Inches(11)
                height = Inches(4)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = content

                # Форматируем текст
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.alignment = PP_ALIGN.LEFT

        except Exception as e:
            logger.warning(f"Не удалось добавить контент: {e}")

        return slide

    def save_to_bytes(self) -> io.BytesIO:
        bytes_io = io.BytesIO()
        self.prs.save(bytes_io)
        bytes_io.seek(0)
        return bytes_io

    def get_slide_count(self) -> int:
        return len(self.prs.slides)
